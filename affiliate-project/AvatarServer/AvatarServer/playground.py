from pptx import Presentation
from pptx.util import Inches

# 创建 Presentation 对象
prs = Presentation()

# 添加标题页
slide_layout = prs.slide_layouts[0]  # 使用标题幻灯片布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "欢迎使用 PPTX 创建工具"
subtitle.text = "这是一个简单的示例"

# 添加内容页
slide_layout = prs.slide_layouts[1]  # 使用标题和内容幻灯片布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "第一页"
content.text = "这是第一页的内容"

# 保存 PPTX 文件
prs.save("example.pptx")
